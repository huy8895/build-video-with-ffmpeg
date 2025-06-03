from pptx import Presentation
import json
import os

pptx_file = "input/slides.pptx"

if not os.path.exists(pptx_file):
    raise FileNotFoundError("slides.pptx not found")

prs = Presentation(pptx_file)
slide_data = []

for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            content = shape.text.strip()
            if content:
                texts.append(content)
    slide_data.append({
        "slide_number": i + 1,
        "text": "\n".join(texts)
    })

os.makedirs("outputs", exist_ok=True)
with open("outputs/slide_content.json", "w", encoding="utf-8") as f:
    json.dump(slide_data, f, indent=2, ensure_ascii=False)

print("✅ Extracted text saved to outputs/slide_content.json")
